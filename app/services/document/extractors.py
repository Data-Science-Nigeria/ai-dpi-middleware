"""Text extractors for office and structured document formats.

Each extractor receives raw file bytes and returns a list of page/slide/sheet
dicts in the same shape as ocr.py: {"page": int, "text": str, "method": str}.
This keeps pipeline.py's _build_chunks() reusable across all formats.
"""

from __future__ import annotations


# ── DOCX ─────────────────────────────────────────────────────────────────────

class _CharPager:
    """Accumulate text fragments and emit a "page" every ``page_chars`` chars.

    Mirrors the buffering the DOCX extractor previously did inline: fragments
    are joined with newlines, and whenever the buffered length reaches the
    threshold a page is flushed with an incrementing page number.
    """

    def __init__(self, page_chars: int):
        self.page_chars = page_chars
        self.pages: list[dict] = []
        self.buf: list[str] = []
        self.buf_len = 0
        self.page_num = 1

    def add(self, text: str) -> None:
        self.buf.append(text)
        self.buf_len += len(text)
        if self.buf_len >= self.page_chars:
            self.flush()

    def flush(self) -> None:
        text = "\n".join(self.buf).strip()
        if text:
            self.pages.append({"page": self.page_num, "text": text, "method": "native"})
            self.page_num += 1
        self.buf = []
        self.buf_len = 0


def _add_docx_paragraphs(doc, pager: _CharPager) -> None:
    for para in doc.paragraphs:
        t = para.text.strip()
        if not t:
            continue
        pager.add(t)


def _add_docx_tables(doc, pager: _CharPager) -> None:
    for table in doc.tables:
        for row in table.rows:
            row_text = " | ".join(c.text.strip() for c in row.cells if c.text.strip())
            if row_text:
                pager.add(row_text)


def extract_docx(file_bytes: bytes) -> list[dict]:
    try:
        import docx  # python-docx
    except ImportError as exc:
        raise RuntimeError("python-docx is required. Install: uv add python-docx") from exc

    import io
    doc = docx.Document(io.BytesIO(file_bytes))

    # Group paragraphs into logical pages (no real page breaks in DOCX without COM).
    # We accumulate text and emit a "page" every PAGE_CHARS characters so that
    # downstream chunking works on reasonably-sized units.
    PAGE_CHARS = 3000
    pager = _CharPager(PAGE_CHARS)

    _add_docx_paragraphs(doc, pager)
    _add_docx_tables(doc, pager)  # Tables

    pager.flush()
    return pager.pages if pager.pages else [{"page": 1, "text": "", "method": "native"}]


# ── PPTX ─────────────────────────────────────────────────────────────────────

def _extract_pptx_body_texts(slide) -> list[str]:
    """Non-title shape text for a slide, in shape/paragraph order."""
    texts: list[str] = []
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        if shape == slide.shapes.title:
            continue  # already added
        for para in shape.text_frame.paragraphs:
            t = para.text.strip()
            if t:
                texts.append(t)
    return texts


def _extract_pptx_slide_parts(slide) -> list[str]:
    parts: list[str] = []

    # Slide title first (if present)
    if slide.shapes.title and slide.shapes.title.text.strip():
        parts.append(slide.shapes.title.text.strip())

    parts.extend(_extract_pptx_body_texts(slide))

    # Speaker notes
    if slide.has_notes_slide:
        notes_tf = slide.notes_slide.notes_text_frame
        notes = notes_tf.text.strip() if notes_tf else ""
        if notes:
            parts.append(f"[Notes] {notes}")

    return parts


def extract_pptx(file_bytes: bytes) -> list[dict]:
    try:
        from pptx import Presentation
        from pptx.util import Pt
    except ImportError as exc:
        raise RuntimeError("python-pptx is required. Install: uv add python-pptx") from exc

    import io
    prs = Presentation(io.BytesIO(file_bytes))
    pages: list[dict] = []

    for slide_num, slide in enumerate(prs.slides, start=1):
        pages.append({
            "page": slide_num,
            "text": "\n".join(_extract_pptx_slide_parts(slide)),
            "method": "native",
        })

    return pages if pages else [{"page": 1, "text": "", "method": "native"}]


# ── XLSX / XLS ────────────────────────────────────────────────────────────────

def extract_xlsx(file_bytes: bytes) -> list[dict]:
    try:
        import openpyxl
    except ImportError as exc:
        raise RuntimeError("openpyxl is required. Install: uv add openpyxl") from exc

    import io
    wb = openpyxl.load_workbook(io.BytesIO(file_bytes), read_only=True, data_only=True)
    pages: list[dict] = []

    for sheet_num, sheet in enumerate(wb.worksheets, start=1):
        rows: list[str] = []
        for row in sheet.iter_rows(values_only=True):
            cells = [str(c) for c in row if c is not None and str(c).strip()]
            if cells:
                rows.append(" | ".join(cells))
        pages.append({
            "page": sheet_num,
            "text": f"Sheet: {sheet.title}\n" + "\n".join(rows),
            "method": "native",
        })

    wb.close()
    return pages if pages else [{"page": 1, "text": "", "method": "native"}]


# ── CSV ───────────────────────────────────────────────────────────────────────

def extract_csv(file_bytes: bytes) -> list[dict]:
    import csv
    import io

    text_io = io.StringIO(file_bytes.decode("utf-8", errors="replace"))
    reader = csv.reader(text_io)
    rows = [" | ".join(cell.strip() for cell in row if cell.strip()) for row in reader]
    rows = [r for r in rows if r]
    return [{"page": 1, "text": "\n".join(rows), "method": "native"}]


# ── RTF ───────────────────────────────────────────────────────────────────────

def extract_rtf(file_bytes: bytes) -> list[dict]:
    try:
        from striprtf.striprtf import rtf_to_text
    except ImportError as exc:
        raise RuntimeError("striprtf is required. Install: uv add striprtf") from exc

    text = rtf_to_text(file_bytes.decode("latin-1", errors="replace"))
    return [{"page": 1, "text": text.strip(), "method": "native"}]


# ── ODT ───────────────────────────────────────────────────────────────────────

def extract_odt(file_bytes: bytes) -> list[dict]:
    try:
        from odf import text as odf_text
        from odf.opendocument import load as odf_load
        from odf.teletype import extractText
    except ImportError as exc:
        raise RuntimeError("odfpy is required. Install: uv add odfpy") from exc

    import io
    doc = odf_load(io.BytesIO(file_bytes))
    paragraphs = doc.getElementsByType(odf_text.P)
    parts = [extractText(p).strip() for p in paragraphs]
    parts = [p for p in parts if p]
    return [{"page": 1, "text": "\n".join(parts), "method": "native"}]


# ── Plain text / Markdown ─────────────────────────────────────────────────────

def extract_plaintext(file_bytes: bytes) -> list[dict]:
    text = file_bytes.decode("utf-8", errors="replace").strip()
    return [{"page": 1, "text": text, "method": "native"}]


# ── Dispatch table ────────────────────────────────────────────────────────────

_EXTRACTORS: dict[str, callable] = {
    ".docx": extract_docx,
    ".doc":  extract_docx,   # python-docx handles older .doc via compatibility
    ".pptx": extract_pptx,
    ".ppt":  extract_pptx,
    ".xlsx": extract_xlsx,
    ".xls":  extract_xlsx,
    ".csv":  extract_csv,
    ".rtf":  extract_rtf,
    ".odt":  extract_odt,
    ".txt":  extract_plaintext,
    ".md":   extract_plaintext,
}

SUPPORTED_EXTENSIONS: frozenset[str] = frozenset(_EXTRACTORS)


def extract(file_bytes: bytes, filename: str) -> list[dict]:
    """Dispatch to the right extractor based on file extension."""
    ext = ("." + filename.rsplit(".", 1)[-1].lower()) if "." in filename else ""
    extractor = _EXTRACTORS.get(ext)
    if extractor is None:
        raise ValueError(
            f"No extractor for extension {ext!r}. "
            f"Supported: {sorted(SUPPORTED_EXTENSIONS)}"
        )
    return extractor(file_bytes)
